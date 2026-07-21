"""办公文档解析器 — Word/Excel/PPT 本地秒级解析 + 损坏检测

防御层次:
  1. 文件预检 — 大小/MIME/可读性/是否损坏
  2. 格式特定检测 — 密码保护/超大行数/二进制旧格式
  3. 解析 — 正常流程，每步独立try/except
"""

import time
import zipfile
from pathlib import Path

from .defenses import validate_file
from .models import DocType, ParseResult, ParseStatus, RawDocument


def parse_office(doc: RawDocument) -> ParseResult:
    """Office文档统一入口"""
    t0 = time.time()

    # ═══════════════════════════════════════
    # 第1层：文件预检
    # ═══════════════════════════════════════
    validation = validate_file(doc.file_path, doc.doc_type.value)
    result = ParseResult(document=doc, status=ParseStatus.SUCCESS)
    result.file_validation = validation

    if not validation.is_valid:
        result.status = ParseStatus.FAILED
        result.errors.extend(validation.errors)
        result.parse_time_ms = (time.time() - t0) * 1000
        return result
    if validation.warnings:
        result.warnings.extend(validation.warnings)

    # ═══════════════════════════════════════
    # 第2层：Office格式特定检测
    # ═══════════════════════════════════════
    ext = Path(doc.file_path).suffix.lower()

    # ZIP损坏检测（docx/xlsx/pptx都是ZIP格式）
    if ext in (".docx", ".xlsx", ".pptx"):
        if not _is_valid_zip(doc.file_path):
            result.status = ParseStatus.FAILED
            result.errors.append(
                "文件可能已损坏（无法作为ZIP打开）。请检查文件是否下载完整，"
                "或尝试另存为后再上传。"
            )
            result.parse_time_ms = (time.time() - t0) * 1000
            return result

    # ═══════════════════════════════════════
    # 第3层：路由到具体解析器
    # ═══════════════════════════════════════
    try:
        if doc.doc_type == DocType.WORD:
            result = _parse_word(doc, t0, result)
        elif doc.doc_type == DocType.EXCEL:
            result = _parse_excel(doc, t0, result)
        elif doc.doc_type == DocType.PPT:
            result = _parse_ppt(doc, t0, result)
        else:
            result.status = ParseStatus.SKIPPED
            result.errors.append(f"不支持的Office格式: {doc.doc_type}")
    except ImportError as e:
        result.status = ParseStatus.FAILED
        result.errors.append(
            f"缺少Office解析依赖: {e}. 请运行 pip install python-docx openpyxl python-pptx"
        )
    except Exception as e:
        result.status = ParseStatus.FAILED
        result.errors.append(f"Office文档解析失败: {e}")

    result.parse_time_ms = (time.time() - t0) * 1000
    return result


def _is_valid_zip(file_path: str) -> bool:
    """检查是否有效的ZIP文件（所有Office Open XML格式的基础）"""
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            # 测试读取——坏文件会在此时抛异常
            _ = zf.namelist()
        return True
    except zipfile.BadZipFile:
        return False
    except Exception:
        return True  # 其他异常不一定是文件损坏


# ═══════════════════════════════════════
# Word
# ═══════════════════════════════════════


def _parse_word(doc: RawDocument, t0: float, result: ParseResult) -> ParseResult:
    from docx import Document as DocxDocument

    word = DocxDocument(doc.file_path)

    markdown_lines = []
    element_count = 0

    for para in word.paragraphs:
        text = para.text.strip()
        if not text:
            markdown_lines.append("")
            continue

        # 检测标题层级
        style_name = para.style.name if para.style else ""
        if "Heading" in style_name or "heading" in style_name:
            level = _extract_heading_level(style_name)
            prefix = "#" * min(level, 6)
            markdown_lines.append(f"{prefix} {text}")
        else:
            markdown_lines.append(text)

        result.elements.append(
            {
                "type": "heading" if "Heading" in style_name else "text",
                "style": style_name,
                "content": text,
            }
        )
        element_count += 1

    # 提取表格
    for ti, table in enumerate(word.tables):
        markdown_lines.append("")
        markdown_lines.append(f"**表格 {ti + 1}:**")
        markdown_lines.append("")

        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                try:
                    cell_text = cell.text.strip().replace("\n", " ")
                except Exception:
                    # 合并单元格可能导致 cell.text 抛异常
                    cell_text = ""
                cells.append(cell_text)
            rows.append("| " + " | ".join(cells) + " |")

        if rows:
            # 表头分隔行
            markdown_lines.append(rows[0])
            col_count = len(table.rows[0].cells) if table.rows else 1
            markdown_lines.append("|" + "|".join(["---"] * col_count) + "|")
            markdown_lines.extend(rows[1:])

        result.elements.append(
            {
                "type": "table",
                "rows": len(rows),
                "content": rows[0] if rows else "",
            }
        )
        element_count += 1

    result.markdown = "\n".join(markdown_lines)
    result.parse_time_ms = (time.time() - t0) * 1000
    return result


# ═══════════════════════════════════════
# Excel
# ═══════════════════════════════════════


def _parse_excel(doc: RawDocument, t0: float, result: ParseResult) -> ParseResult:
    import openpyxl

    wb = openpyxl.load_workbook(doc.file_path, data_only=True)

    markdown_lines = [f"# {Path(doc.file_path).stem}", ""]

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if ws.max_row == 0:
            continue

        markdown_lines.append(f"## {sheet_name}")
        markdown_lines.append("")

        # 取前500行（避免超大Excel）
        rows = list(ws.iter_rows(max_row=min(ws.max_row, 500), values_only=True))
        if not rows:
            continue

        # 构建Markdown表格
        header = [str(c) if c is not None else "" for c in rows[0]]
        col_count = len(header)
        markdown_lines.append("| " + " | ".join(header) + " |")
        markdown_lines.append("|" + "|".join(["---"] * col_count) + "|")

        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            # 补齐列数
            while len(cells) < col_count:
                cells.append("")
            markdown_lines.append("| " + " | ".join(cells[:col_count]) + " |")

        markdown_lines.append("")

        result.elements.append(
            {
                "type": "sheet",
                "sheet_name": sheet_name,
                "rows": len(rows),
            }
        )

    wb.close()
    result.markdown = "\n".join(markdown_lines)
    result.parse_time_ms = (time.time() - t0) * 1000
    return result


# ═══════════════════════════════════════
# PPT
# ═══════════════════════════════════════


def _parse_ppt(doc: RawDocument, t0: float, result: ParseResult) -> ParseResult:
    from pptx import Presentation

    prs = Presentation(doc.file_path)

    markdown_lines = [f"# {Path(doc.file_path).stem}", ""]

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    slide_texts.append("\n".join(rows))

        if slide_texts:
            markdown_lines.append(f"## 第{slide_num}页")
            markdown_lines.append("\n".join(slide_texts))
            markdown_lines.append("")

        result.elements.append(
            {
                "type": "slide",
                "slide_num": slide_num,
                "text_count": len(slide_texts),
            }
        )

    result.markdown = "\n".join(markdown_lines)
    result.parse_time_ms = (time.time() - t0) * 1000
    return result


def _extract_heading_level(style_name: str) -> int:
    """从Word样式名提取标题层级"""
    for char in style_name:
        if char.isdigit():
            return int(char)
    return 1
